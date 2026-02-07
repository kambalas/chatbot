from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def extract_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
    return "\n".join(parts)
